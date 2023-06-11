from flask import Flask, render_template, request, send_file
import io
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx import Presentation
import os
import openai
import config
openai.api_key = config.OPENAI_API_KEY

app = Flask(__name__)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    input_name = request.form['input_name']

    # Generate PPT file
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = input_name
    subtitle.text = "Generated using GPT"

    for i in range(9):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = 'Slide %d' % (i + 2)
        tf = body_shape.text_frame
        tf.text = 'This is bullet slide %d' % (i + 2)
        p = tf.add_paragraph()
        p.text = 'Bullet 1'
        p.level = 1
        p = tf.add_paragraph()
        p.text = 'Bullet 2'
        p.level = 2

    # Save PPT file to memory stream
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)

    return send_file(ppt_stream, attachment_filename=f"{input_name}.pptx", as_attachment=True)
