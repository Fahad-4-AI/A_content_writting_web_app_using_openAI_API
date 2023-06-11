from flask import Flask, render_template, request, send_file,request , send_from_directory
from pptx import Presentation
import os
import openai
import config
openai.api_key = config.OPENAI_API_KEY

app = Flask(__name__)


@app.route("/")
def home():
    return render_template("index.html")

@app.route('/index')
def index():
    return render_template("index.html")




# Define the route to display definitions
@app.route('/definitions', methods=['GET', 'POST'])
def definitions():
    if request.method == 'POST':
        term = request.form['term']
        definition = generate_definition(term)
        return render_template('definitions.html', term=term, definition=definition)
    return render_template('definitions.html')

# Function to generate definitions using GPT API
def generate_definition(term):
    prompt = f"Define {term}?"
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=prompt,
        temperature=0.7,
        max_tokens=256,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0,
    )
    definition = response.choices[0].text.strip()
    return definition





@app.route('/explanations', methods=['GET', 'POST'])
def explanations():
    if request.method == 'POST':
        topic = request.form['topic']
        explanation = generate_explanation(topic)
        return render_template('explanations.html', topic=topic, explanation=explanation)
    return render_template('explanations.html')

def generate_explanation(topic):
    prompt = f"Explain {topic}."
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=prompt,
        temperature=0.7,
        max_tokens=256,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0,)
    explanation = response['choices'][0]['text'].replace('/n','<br>')
    return explanation






# Define the route to display article
@app.route("/article", methods=["GET", "POST"])
def article():
    generated_article = ""
    if request.method == "POST":
        topic = request.form.get("topic")
        article = request.form.get("article")
        response = openai.Completion.create(
            engine="text-davinci-002",
            prompt=f"Write an article about {topic}.\n\n{article}\n\n",
            max_tokens=4000,
            n=1,
            stop=None,
            temperature=0.5,
        )
        generated_article = response.choices[0].text.strip()
    return render_template("article.html", generated_article=generated_article)










# Define route for generating a story
@app.route("/story", methods=["GET", "POST"])
def story():
    generated_story = None
    if request.method == "POST":
        characters = request.form["characters"]
        setting = request.form["setting"]
        plot = request.form["plot"]
        
        # Concatenate characters, setting, and plot into a prompt
        prompt = f"{characters} are in {setting}. {plot}"
        
        # Set up OpenAI API parameters
        completions = openai.Completion.create(
            engine="text-davinci-002",
            prompt=prompt,
            max_tokens=1024,
            n=1,
            stop=None,
            temperature=0.7,
        )

        # Extract the generated story from the API response
        message = completions.choices[0].text
        generated_story = message.strip()
        
    return render_template("story.html", generated_story=generated_story)











# Define route for generating an essay
@app.route("/essay", methods=["GET", "POST"])
def essay():
    generated_essay = None
    if request.method == "POST":
        topic = request.form["topic"]
        prompt = f"Write an essay about {topic}."
        
        # Set up OpenAI API parameters
        completions = openai.Completion.create(
            engine="text-davinci-002",
            prompt=prompt,
            max_tokens=1024,
            n=1,
            stop=None,
            temperature=0.7,
        )

        # Extract the generated essay from the API response
        message = completions.choices[0].text
        generated_essay = message.strip()
        
    return render_template("essay.html", generated_essay=generated_essay)









@app.route("/cpp_programming", methods=["GET", "POST"])
def cpp_programming():
    if request.method == "POST":
        problem = request.form["problem"]
        response = openai.Completion.create(
            engine="text-davinci-002",
            prompt=f"Write a C++ program for: {problem}",
            temperature=0.5,
            max_tokens=1024,
            n=1,
            stop=None,
            frequency_penalty=0,
            presence_penalty=0
        )

        return render_template("cpp_programming.html", response=response.choices[0].text)

    return render_template("cpp_programming.html")









@app.route("/python_programming", methods=["GET", "POST"])
def python_programming():
    if request.method == "POST":
        problem = request.form["problem"]
        response = openai.Completion.create(
            engine="text-davinci-002",
            prompt=f"Write a Python program for: {problem}",
            temperature=0.5,
            max_tokens=1024,
            n=1,
            stop=None,
            frequency_penalty=0,
            presence_penalty=0
        )

        return render_template("python_programming.html", response=response.choices[0].text)

    return render_template("python_programming.html")










""""
from flask import Flask, render_template, request, make_response
import io
import random
from pptx import Presentation
from pptx.util import Inches
import openai



@app.route('/generate_ppt', methods=['GET', 'POST'])
def generate_ppt():
    if request.method == 'POST':
        input_text = request.form['input_text']
        try:
            model_engine = "text-davinci-002"
            prompt = (f"Generate a 10-slide PowerPoint presentation on {input_text}. "
                      "Include bullet points and images to support the presentation. "
                      "Limit the number of bullet points to 5 per slide.")
            completions = openai.Completion.create(
                engine=model_engine,
                prompt=prompt,
                max_tokens=1024,
                n=1,
                stop=None,
                temperature=0.7,
            )

            message = completions.choices[0].text.strip()
            presentation = Presentation()

            for i in range(0,10):
                slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                slide.shapes.title.text = f"Slide {i+1}"
                tf = slide.shapes.placeholders[1].text_frame
                tf.text = message.split("Slide {}".format(i+1))[0].split("Slide {}".format(i))[1]

                bullet_point = ""
                bullet_points = message.split("Slide {}".format(i+1))[1].split("Slide {}".format(i+2))[0]
                bullet_point_list = bullet_points.split("\n")

                for point in bullet_point_list:
                    if point != "":
                        if bullet_point == "":
                            bullet_point = point
                        else:
                            bullet_point += "\n" + point

                p = tf.add_paragraph()
                p.text = bullet_point

                if random.choice([True, False]):
                    image = io.BytesIO(openai.Image.create(completions.choices[0].text).get_bytes())
                    left = top = Inches(1)
                    pic = slide.shapes.add_picture(image, left, top)

            pptx_output = io.BytesIO()
            presentation.save(pptx_output)
            pptx_output.seek(0)

            response = make_response(pptx_output.read())
            response.headers.set('Content-Disposition', 'attachment', filename='output.pptx')
            response.headers.set('Content-Type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
            return response

        except Exception as e:
            print(str(e))
            return "Error: " + str(e)

    return render_template('index.html')"""

         
 
import pptx
import presentation

# Define the route to display generate_ppt
@app.route("/generate_ppt", methods=["GET", "POST"])
def generate_ppt():
    generated_generate_ppt = ""
    if request.method == "POST":
        topic = request.form.get("topic")
        generate_ppt = request.form.get("generate_ppt")
        response = openai.Completion.create(
            engine="text-davinci-002",
            prompt=f"give me 3 pptx slides text data about {topic}.Each slides with a title and explination\n\n{generate_ppt}\n\n",
            max_tokens=2000,
            n=1,
            stop=None,
            temperature=0.5,
        )
# Create a presentation
    presentation = pptx.Presentation()

# Add slides to the presentation
    for slide in response.choices[0].text.split("\n\n"):
     slide_title = slide.split(":")[0]
     slide_content = slide.split(":")[1]
     presentation.slides.add_slide(pptx.Slide(slide_title, slide_content))

# Save the presentation
    presentation.save("presentation.pptx")
    return send_file("presentation.pptx", as_attachment=True)


















#pip install python-pptx






@app.route("/ask_anything")
def ask_anything():
    return render_template("ask_anything.html")

@app.route("/get_answer", methods=["POST"])
def get_answer():
    question = request.form["question"]
    response = openai.Completion.create(
        engine="davinci",
        prompt=f"{question}",
        max_tokens=100,
        n=1,
        stop=None,
        temperature=0.7
    )
    answer = response.choices[0].text
    return render_template("ask_anything.html", answer=answer)

if __name__ == "__main__":
    app.run(debug=True)

       
